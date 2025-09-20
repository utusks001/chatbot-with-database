# app-langchain.py

import streamlit as st
import pandas as pd
import plotly.express as px
import os, tempfile, requests
from pptx import Presentation
from langchain.chains import LLMChain
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_groq import ChatGroq
from langchain.prompts import ChatPromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader

# =====================
# Session State Init
# =====================
if "dfs" not in st.session_state:
    st.session_state.dfs = {}
if "vectorstore" not in st.session_state:
    st.session_state.vectorstore = None
if "uploaded_rag_files" not in st.session_state:
    st.session_state.uploaded_rag_files = []

# =====================
# LLM SETUP
# =====================
def load_llm():
    try:
        return ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.3)
    except Exception:
        return ChatGroq(model="llama-3.3-70b-versatile", temperature=0.3)

llm = load_llm()

# =====================
# Data Analysis Helpers
# =====================
def df_info_text(df: pd.DataFrame) -> str:
    info = f"Rows: {df.shape[0]}, Columns: {df.shape[1]}\n"
    info += "Columns:\n" + ", ".join(df.columns[:30])
    if df.shape[1] > 30:
        info += " ..."
    return info

def detect_column_types(df: pd.DataFrame):
    numeric_cols = df.select_dtypes(include=["int64","float64"]).columns.tolist()
    categorical_cols = df.select_dtypes(include=["object","category"]).columns.tolist()
    datetime_cols = df.select_dtypes(include=["datetime64[ns]"]).columns.tolist()
    return numeric_cols, categorical_cols, datetime_cols

def safe_describe(df):
    try:
        return df.describe(include="all")
    except Exception:
        return pd.DataFrame()

def generate_dataset_insight(df: pd.DataFrame, question: str = None):
    stats = safe_describe(df).reset_index().to_string()
    question_section = f"Question: {question}" if question else ""
    prompt_template = f"""
    You are a data analyst. Based on the dataset stats below:
    {stats}
    {question_section}
    Provide a clear, concise, and natural insight or answer.
    If answer is not available, say: "Answer not available in the provided context."
    """
    prompt = ChatPromptTemplate.from_template(prompt_template)
    chain = prompt | llm
    return chain.invoke({}).content

# =====================
# OCR Helper
# =====================
OCR_SPACE_API_KEY = st.secrets.get("OCR_SPACE_API_KEY", "")

def ocr_space(file_path):
    """OCR using OCR.Space API"""
    with open(file_path, "rb") as f:
        r = requests.post(
            "https://api.ocr.space/parse/image",
            files={"filename": f},
            data={"apikey": OCR_SPACE_API_KEY, "language":"eng"}
        )
    try:
        result = r.json()
        text = ""
        for parsed in result.get("ParsedResults", []):
            text += parsed.get("ParsedText","") + "\n"
        return text
    except:
        return ""

# =====================
# Document Loaders (without NLTK)
# =====================
def load_pptx(file_path):
    prs = Presentation(file_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return [Document(page_content=text)]

def load_document(file_path, file_type):
    file_type = file_type.lower()
    if file_type == ".pdf":
        return PyPDFLoader(file_path).load()
    elif file_type == ".txt":
        return TextLoader(file_path).load()
    elif file_type == ".docx":
        return Docx2txtLoader(file_path).load()
    elif file_type in [".pptx", ".ppt"]:
        return load_pptx(file_path)
    elif file_type in [".jpg",".jpeg",".png",".bmp",".gif"]:
        text = ocr_space(file_path)
        if text.strip():
            return [Document(page_content=text)]
        else:
            return []
    else:
        return []

def process_rag_files(uploaded_files):
    docs = []
    for uploaded_file in uploaded_files:
        suffix = os.path.splitext(uploaded_file.name)[1].lower()
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(uploaded_file.read())
            tmp_path = tmp.name
        docs.extend(load_document(tmp_path, suffix))
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    return splitter.split_documents(docs)

# =====================
# UI
# =====================
st.set_page_config(page_title="ü§ñüìä Data & Document Chatbot", layout="wide")
st.title("ü§ñüìä Chatbot Dashboard: Data Analysis & Advanced RAG")

tab1, tab2 = st.tabs(["üìà Data Analysis", "üìö RAG Advanced"])

# ====== Data Analysis ======
with tab1:
    uploaded_file = st.file_uploader("Upload CSV/Excel for analysis", type=["csv","xls","xlsx"])
    df = None
    if uploaded_file:
        if uploaded_file.name.endswith(".csv"):
            df = pd.read_csv(uploaded_file)
        else:
            xls = pd.ExcelFile(uploaded_file)
            st.session_state.dfs = {sheet: pd.read_excel(uploaded_file, sheet_name=sheet) for sheet in xls.sheet_names}
            selected_sheets = st.multiselect("üìë Select Sheet", list(st.session_state.dfs.keys()), default=list(st.session_state.dfs.keys())[:1])
            if selected_sheets:
                df_list = []
                for s in selected_sheets:
                    temp = st.session_state.dfs[s].copy()
                    temp["SheetName"] = s
                    df_list.append(temp)
                df = pd.concat(df_list, ignore_index=True)
    if df is not None:
        st.dataframe(df.head(10))
        numeric_cols, categorical_cols, datetime_cols = detect_column_types(df)
        st.write(f"Numeric: {numeric_cols}")
        st.write(f"Categorical: {categorical_cols}")
        st.write(f"Datetime: {datetime_cols}")
        st.text(df_info_text(df))
        st.write(f"**Data shape:** {df.shape}")        

        st.subheader("‚öôÔ∏è Select Columns for Visualization")
        x_axis = st.selectbox("X Axis", df.columns)
        y_axis = st.selectbox("Y Axis", df.columns)
        if x_axis and y_axis:
            x_is_num = x_axis in numeric_cols or x_axis in datetime_cols
            y_is_num = y_axis in numeric_cols
            x_is_cat = x_axis in categorical_cols
            y_is_cat = y_axis in categorical_cols
            fig = None
            if x_is_num and y_is_num:
                fig = px.scatter(df, x=x_axis, y=y_axis, title=f"üìà Scatter {y_axis} vs {x_axis}")
            elif x_axis in datetime_cols and y_is_num:
                fig = px.line(df, x=x_axis, y=y_axis, title=f"üìà Trend {y_axis} vs {x_axis}")
            elif x_is_cat and y_is_num:
                fig = px.bar(df, x=x_axis, y=y_axis, title=f"üìä Bar {y_axis} per {x_axis}")
            elif x_is_num and y_is_cat:
                fig = px.bar(df, x=y_axis, y=x_axis, title=f"üìä Bar {x_axis} per {y_axis}")
            elif x_is_cat and y_is_cat:
                crosstab = pd.crosstab(df[x_axis], df[y_axis])
                fig = px.imshow(crosstab, title=f"üî¢ Frequency {x_axis} vs {y_axis}")
            if fig:
                st.plotly_chart(fig, use_container_width=True)

        st.subheader("üí¨ Chatbot Data Analysis")
        q = st.text_input("Ask about dataset")
        if q:
            st.write(generate_dataset_insight(df, question=q))

# ====== RAG Advanced ======
with tab2:
    uploaded_files = st.file_uploader(
        "Upload PDF/DOCX/PPTX/TXT/Image (OCR)",
        type=["pdf","docx","pptx","txt","jpg","jpeg","png","bmp","gif"],
        accept_multiple_files=True
    )
    if uploaded_files:
        st.session_state.uploaded_rag_files = uploaded_files
        st.write("üìÑ Uploaded files:")
        for f in uploaded_files:
            st.write(f.name)

    if st.session_state.uploaded_rag_files:
        if st.button("üöÄ Process Files to Vectorstore"):
            with st.spinner("üì¶ Processing documents..."):
                docs = process_rag_files(st.session_state.uploaded_rag_files)
                if docs:
                    embeddings = HuggingFaceEmbeddings()
                    st.session_state.vectorstore = FAISS.from_documents(docs, embeddings)
                    st.success(f"‚úÖ Documents processed: {len(docs)} chunks")
                else:
                    st.error("No valid documents to process")

        st.subheader("üí¨ Chatbot RAG Query")
        user_query = st.text_input("Ask a question about uploaded documents")
        if user_query and st.session_state.vectorstore:
            retriever = st.session_state.vectorstore.as_retriever(search_kwargs={"k":3})
            related_docs = retriever.get_relevant_documents(user_query)
            context_text = "\n\n".join([doc.page_content for doc in related_docs])
            prompt_template = f"""
            You are a helpful AI assistant. Answer the user question based on the following documents:
            {context_text}
            Question: {user_query}
            """
            prompt = ChatPromptTemplate.from_template(prompt_template)
            chain = prompt | llm
            answer = chain.invoke({}).content
            st.write(answer)
