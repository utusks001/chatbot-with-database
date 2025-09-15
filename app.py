import os
import requests
from io import BytesIO
import streamlit as st
from dotenv import load_dotenv

# File parsing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from PIL import Image

# LangChain / VectorStore / Embeddings / LLM
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# LangGraph + LangSmith
from typing import TypedDict
from langchain_core.prompts import ChatPromptTemplate
from langgraph.graph import StateGraph, END
from langsmith import Client

# -------------------------
# Config
# -------------------------
load_dotenv()
OCR_SPACE_API_KEY = os.getenv("OCR_SPACE_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "")
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
LANGCHAIN_API_KEY = os.getenv("LANGCHAIN_API_KEY", "")

# Aktifkan LangSmith tracing (jika API Key ada)
if LANGCHAIN_API_KEY:
    os.environ["LANGCHAIN_TRACING_V2"] = "true"
    os.environ["LANGCHAIN_API_KEY"] = LANGCHAIN_API_KEY
    client = Client()

st.set_page_config(
    page_title="Gemini + Groq Chatbot — LangGraph + LangSmith",
    page_icon="🤖",
    layout="wide"
)

# -------------------------
# Embeddings + Splitter
# -------------------------
EMBEDDINGS = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
SPLITTER = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)

# -------------------------
# File extractors
# -------------------------
def extract_text_from_pdf(file_bytes: BytesIO):
    text = ""
    try:
        reader = PdfReader(file_bytes)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PDF: {e}")
    return text

def extract_text_from_txt(file_bytes: BytesIO):
    try:
        return file_bytes.read().decode("utf-8", errors="ignore")
    except Exception as e:
        st.warning(f"⚠️ Gagal baca TXT: {e}")
        return ""

def extract_text_from_docx(file_bytes: BytesIO):
    text = ""
    try:
        file_bytes.seek(0)
        doc = DocxDocument(file_bytes)
        for p in doc.paragraphs:
            if p.text:
                text += p.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak DOCX: {e}")
    return text

def extract_text_from_pptx(file_bytes: BytesIO):
    text = ""
    try:
        file_bytes.seek(0)
        prs = PptxPresentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PPTX: {e}")
    return text

def extract_text_from_image(file_bytes: BytesIO, filename="upload.png"):
    if not OCR_SPACE_API_KEY:
        st.warning("⚠️ OCR_SPACE_API_KEY tidak ditemukan di .env")
        return ""
    try:
        file_bytes.seek(0)
        response = requests.post(
            "https://api.ocr.space/parse/image",
            files={"file": (filename, file_bytes, "image/png")},
            data={"apikey": OCR_SPACE_API_KEY, "language": "eng"},
        )
        result = response.json()
        if result.get("IsErroredOnProcessing"):
            st.warning("⚠️ OCR.Space gagal: " + str(result.get("ErrorMessage", ['Unknown error'])))
            return ""
        text = "\n".join([p["ParsedText"] for p in result.get("ParsedResults", []) if "ParsedText" in p])
        return text.strip()
    except Exception as e:
        st.warning(f"⚠️ OCR.Space error: {e}")
        return ""

def extract_text_from_file(uploaded_file):
    name = uploaded_file.name.lower()
    raw = uploaded_file.read()
    bio = BytesIO(raw)

    if name.endswith(".pdf"):
        return extract_text_from_pdf(bio)
    elif name.endswith(".txt"):
        return extract_text_from_txt(BytesIO(raw))
    elif name.endswith(".docx"):
        return extract_text_from_docx(BytesIO(raw))
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(BytesIO(raw))
    elif name.endswith((".jpg", ".jpeg", ".png", ".gif", ".bmp", ".jfif")):
        return extract_text_from_image(BytesIO(raw), filename=uploaded_file.name)
    else:
        st.warning(f"⚠️ Tipe file `{uploaded_file.name}` tidak didukung.")
        return ""

# -------------------------
# Build documents & FAISS
# -------------------------
def build_documents_from_uploads(uploaded_files):
    docs = []
    for f in uploaded_files:
        text = extract_text_from_file(f)
        if not text or not text.strip():
            continue
        chunks = SPLITTER.split_text(text)
        for i, chunk in enumerate(chunks):
            docs.append(Document(page_content=chunk, metadata={"source_file": f.name, "chunk_id": i}))
    return docs

def build_faiss_from_documents(docs):
    if not docs:
        return None
    vs = FAISS.from_documents(docs, embedding=EMBEDDINGS)
    return vs

# -------------------------
# LangGraph State
# -------------------------
class GraphState(TypedDict):
    question: str
    context: str
    answer: str

prompt_template = ChatPromptTemplate.from_messages([
    ("system", "Jawablah secara akurat, jelas, dan detil berdasarkan konteks berikut. "
               "Jika jawaban tidak ada, katakan 'Jawaban tidak tersedia dalam konteks'.\n\nKonteks:\n{context}"),
    ("human", "{question}")
])

# Node: retriever
def retrieve_node(state: GraphState):
    results = st.session_state.vector_store.similarity_search(state["question"], k=5)
    ctx = "\n\n".join([doc.page_content for doc in results])
    return {"question": state["question"], "context": ctx, "answer": ""}

# Node: LLM
def llm_node(state: GraphState):
    if st.session_state.model_choice.startswith("Gemini"):
        from langchain_google_genai import ChatGoogleGenerativeAI
        llm = ChatGoogleGenerativeAI(
            model="gemini-2.5-flash", temperature=0.2, google_api_key=GOOGLE_API_KEY
        )
    else:
        from langchain_groq import ChatGroq
        llm = ChatGroq(
            model_name="llama-3.3-70b-versatile", temperature=0.2, groq_api_key=GROQ_API_KEY
        )
    chain = prompt_template | llm
    response = chain.invoke({"context": state["context"], "question": state["question"]})
    return {"answer": response.content, "question": state["question"], "context": state["context"]}

# Bangun graph
workflow = StateGraph(GraphState)
workflow.add_node("retriever", retrieve_node)
workflow.add_node("llm", llm_node)
workflow.set_entry_point("retriever")
workflow.add_edge("retriever", "llm")
workflow.add_edge("llm", END)
graph = workflow.compile()

# -------------------------
# Streamlit UI
# -------------------------
st.title("🤖 Gemini 2.5 Flash + Groq — LangGraph + LangSmith")
st.write("Upload banyak file (PDF, TXT, DOCX, PPTX, Images). OCR untuk gambar via OCR.Space.")

# Sidebar upload
st.sidebar.header("📂 Upload & Build")
uploaded_files = st.sidebar.file_uploader(
    "Upload files", type=["pdf", "txt", "docx", "pptx", "jpg", "jpeg", "png", "gif", "bmp", "jfif"],
    accept_multiple_files=True
)
build_btn = st.sidebar.button("🚀 Build Vector Store")
clear_btn = st.sidebar.button("🧹 Reset vector store")

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "indexed_files" not in st.session_state:
    st.session_state.indexed_files = []
if "model_choice" not in st.session_state:
    st.session_state.model_choice = "Gemini 2.5 Flash (Google)"

if clear_btn:
    st.session_state.vector_store = None
    st.session_state.indexed_files = []
    st.success("Vector store di-reset.")

if build_btn:
    if not uploaded_files:
        st.sidebar.warning("Silakan upload minimal 1 file.")
    else:
        with st.spinner("📦 Memproses file..."):
            docs = build_documents_from_uploads(uploaded_files)
            if not docs:
                st.sidebar.error("Tidak ada teks berhasil diekstrak.")
            else:
                vs = build_faiss_from_documents(docs)
                st.session_state.vector_store = vs
                st.session_state.indexed_files = [f.name for f in uploaded_files]
                st.sidebar.success(f"Vector store siap. Dokumen: {len(st.session_state.indexed_files)} | Chunk: {len(docs)}")

if st.session_state.indexed_files:
    st.markdown("**Dokumen terindeks:**")
    st.write(" • " + "\n • ".join(st.session_state.indexed_files))

# Pilih LLM Provider
st.session_state.model_choice = st.sidebar.radio(
    "Pilih LLM Provider:", ["Gemini 2.5 Flash (Google)", "Groq (llama-3.3-70b-versatile)"],
    index=0
)

# Query
prompt = st.text_input("Tanyakan sesuatu berdasarkan dokumen:")
ask_btn = st.button("Tanyakan")

if ask_btn:
    if not prompt.strip():
        st.warning("Masukkan pertanyaan terlebih dahulu.")
    elif st.session_state.vector_store is None:
        st.info("Belum ada vector store. Upload & build terlebih dahulu.")
    else:
        with st.spinner("🚀 Menjalankan LangGraph..."):
            result = graph.invoke({"question": prompt, "context": "", "answer": ""})
        st.subheader("💬 Jawaban")
        st.write(result["answer"])
