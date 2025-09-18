# app-multisheet-rag.py

# app-multisheet-rag.py

import streamlit as st
import os, io
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns

from pathlib import Path
import toml
from dotenv import load_dotenv, set_key

# LangChain & RAG
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain.docstore.document import Document
from langchain.chains import RetrievalQA
from langchain_experimental.tools.python.tool import PythonREPLTool
from langchain.agents import initialize_agent, Tool, AgentType

# File loaders
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import easyocr
from PIL import Image

# Load dotenv kalau lokal
dotenv_path = Path(".env")
if dotenv_path.exists():
    load_dotenv(dotenv_path)

# ====== Helper deteksi local vs Streamlit Cloud ======
def is_streamlit_cloud():
    return os.environ.get("STREAMLIT_RUNTIME") is not None

# ====== Sidebar API Key ======
with st.sidebar:
    st.header("🔑 Konfigurasi API Key")

    GOOGLE_API_KEY = (
        st.secrets.get("GOOGLE_API_KEY", "")
        or os.getenv("GOOGLE_API_KEY", "")
    )

    if not GOOGLE_API_KEY:
        GOOGLE_API_KEY = st.text_input(
            "Masukkan GOOGLE_API_KEY (buat di https://aistudio.google.com/apikey)",
            type="password"
        )
        if GOOGLE_API_KEY:
            os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
            st.session_state["GOOGLE_API_KEY"] = GOOGLE_API_KEY
            st.success("✅ GOOGLE_API_KEY berhasil dimasukkan")

            if is_streamlit_cloud():
                st.info("ℹ️ Running di Streamlit Cloud → gunakan Settings → Secrets untuk simpan permanen")
            else:
                set_key(dotenv_path, "GOOGLE_API_KEY", GOOGLE_API_KEY)
                st.success("✅ API Key disimpan ke .env (lokal)")
    else:
        st.success("✅ GOOGLE_API_KEY berhasil dimuat")

    # Embeddings toggle
    st.subheader("⚙️ Embeddings Settings")
    use_hf_embeddings = st.checkbox("Pakai HuggingFace embeddings saja", value=False)
    st.session_state["USE_HF_EMBEDDINGS"] = use_hf_embeddings

# ====== Helper Functions ======
def safe_describe(df: pd.DataFrame):
    try:
        return df.describe(include="all").transpose()
    except Exception as e:
        return pd.DataFrame({"Error": [str(e)]})

def df_info_text(df: pd.DataFrame):
    buf = io.StringIO()
    df.info(buf=buf)
    return buf.getvalue()

# ====== Helper: Deteksi kolom numerik vs kategorikal ======
def detect_data_types(df: pd.DataFrame):
    categorical_cols = df.select_dtypes(include=["object", "category"]).columns.tolist()
    numeric_cols = df.select_dtypes(include=["number"]).columns.tolist()
    return categorical_cols, numeric_cols

# ====== OCR Helper (EasyOCR) ======
ocr_reader = easyocr.Reader(["en"], gpu=False)

def ocr_image(file):
    img = Image.open(file)
    results = ocr_reader.readtext(img)
    text = "\n".join([res[1] for res in results])
    return text

# ====== Build Vectorstore ======
def build_vectorstore(files):
    """Bangun vectorstore dari dokumen upload (multi-file)."""
    docs = []

    for file in files:
        name = file.name

        if name.endswith(".txt"):
            text = file.read().decode("utf-8")
            docs.append(Document(page_content=text, metadata={"source": name}))

        elif name.endswith(".pdf"):
            reader = PdfReader(file)
            text = "\n".join([page.extract_text() or "" for page in reader.pages])
            docs.append(Document(page_content=text, metadata={"source": name}))

        elif name.endswith(".csv"):
            df = pd.read_csv(file)
            text = df.to_csv(index=False)
            docs.append(Document(page_content=text, metadata={"source": name}))

        elif name.endswith(".xlsx") or name.endswith(".xls"):
            xls = pd.ExcelFile(file)
            for sheet in xls.sheet_names:
                df = pd.read_excel(file, sheet_name=sheet)
                text = f"[{sheet}]\n" + df.to_csv(index=False)
                docs.append(Document(page_content=text, metadata={"source": f"{name}:{sheet}"}))

        elif name.endswith(".docx"):
            doc = docx.Document(file)
            text = "\n".join([para.text for para in doc.paragraphs])
            docs.append(Document(page_content=text, metadata={"source": name}))

        elif name.endswith(".pptx"):
            prs = Presentation(file)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            docs.append(Document(page_content=text, metadata={"source": name}))

        elif name.lower().endswith((".png", ".jpg", ".jpeg", ".bmp")):
            text = ocr_image(file)
            docs.append(Document(page_content=text, metadata={"source": name}))

        else:
            try:
                text = file.read().decode("utf-8", errors="ignore")
                docs.append(Document(page_content=text, metadata={"source": name}))
            except Exception:
                pass

    # Split dokumen
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    split_docs = splitter.split_documents(docs)

    # ====== Embeddings Selection ======
    if st.session_state.get("USE_HF_EMBEDDINGS", False):
        st.info("ℹ️ Mode HuggingFace embeddings dipilih (manual).")
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        return FAISS.from_documents(split_docs, embeddings)

    try:
        embeddings = GoogleGenerativeAIEmbeddings(
            model="models/embedding-001",
            google_api_key=os.getenv("GOOGLE_API_KEY", "")
        )
        return FAISS.from_documents(split_docs, embeddings)
    except Exception as e:
        st.warning(f"⚠️ Gagal pakai Gemini embeddings ({e}). Fallback ke HuggingFace.")
        st.session_state["USE_HF_EMBEDDINGS"] = True
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        return FAISS.from_documents(split_docs, embeddings)


# ====== Main App ======
st.set_page_config(page_title="Data & Document RAG Chatbot", layout="wide")
st.title("📊🤖 Chatbot Analisis Data & Dokumen (RAG + Gemini/HF)")

menu = st.sidebar.radio("Pilih Mode", ["📊 Data Analysis", "📑 RAG Advanced"])

# ========== MODE 1: Data Analysis ==========
if menu == "📊 Data Analysis":
    uploaded_file = st.file_uploader(
        "Upload file Excel/CSV untuk analisa data",
        type=["csv", "xls", "xlsx"]
    )

    if uploaded_file is not None:
        if uploaded_file.name.endswith(".csv"):
            df = pd.read_csv(uploaded_file)
            st.session_state.dfs = {"CSV": df}
        else:
            xls = pd.ExcelFile(uploaded_file)
            st.session_state.dfs = {sheet: pd.read_excel(uploaded_file, sheet_name=sheet) for sheet in xls.sheet_names}

        with st.sidebar:
            st.subheader("📑 Pilih Sheet")
            sheet_names = list(st.session_state.dfs.keys())
            selected_sheets = st.multiselect("Sheet Aktif", sheet_names, default=sheet_names[:1])

        if not selected_sheets:
            st.warning("Pilih minimal satu sheet untuk analisis.")
            st.stop()

        if len(selected_sheets) == 1:
            df = st.session_state.dfs[selected_sheets[0]]
            sheet_label = selected_sheets[0]
        else:
            df_list = []
            for s in selected_sheets:
                temp = st.session_state.dfs[s].copy()
                temp["SheetName"] = s
                df_list.append(temp)
            df = pd.concat(df_list, ignore_index=True)
            sheet_label = ", ".join(selected_sheets)

        st.markdown(f"### 📄 Analisa: {uploaded_file.name} — Sheet(s): {sheet_label}")
        
        # Preview data
        st.write("**Head (10):**")
        st.dataframe(df.head(10))
        st.write("**Tail (10):**")
        st.dataframe(df.tail(10))
        
        # Ringkasan kolom
        categorical_cols, numeric_cols = detect_data_types(df)
        st.write("**Ringkasan Kolom**")
        st.write(f"Kolom Numerik: {numeric_cols}")
        st.write(f"Kolom Kategorikal: {categorical_cols}")
        
        st.write("**Info():**")
        st.text(df_info_text(df))
        
        # Display the shape of the data
        st.write(f"**Data shape:** {df.shape}")
        st.write("                                             ")  
        
        # Display data information
        st.write("**Data information:**")
        for index, (col, dtype) in enumerate(zip(df.columns, df.dtypes)):
            non_null_count = df[col].count()
            st.write(f"{index} | {col}   | {non_null_count} non-null  |  {dtype}") 
        
        # Check missing values again
        missing_values = df.isnull().sum()
        st.write("Missing Values :")
        st.write(missing_values)
        st.write("                                             ")  
        
        # Display the number of duplicates removed
        duplicates_count = df.duplicated().sum()
        st.write(f"Number of Duplicates : {duplicates_count}")
        st.write("                                             ")  
                
        # Remove duplicates
        duplicates_removed = df.drop_duplicates(inplace=True)
        st.write(f"Number of Duplicates Removed: {duplicates_removed}")
        st.write("                                             ")  
        
        # Display Describe
        st.write("**Describe():**")
        st.dataframe(safe_describe(df))
        
        # Display summary statistics of the DataFrame
        st.write("**Summary Statistics:**")
        st.write(df.describe(include="all"))

        # Correlation heatmap
        num_df = df.select_dtypes(include="number")
        if not num_df.empty:
            st.write("**Correlation Heatmap**")
            fig, ax = plt.subplots(figsize=(6, 4))
            sns.heatmap(num_df.corr(), annot=True, cmap="coolwarm", ax=ax)
            st.pyplot(fig)

        # ====== Chatbot Data Analysis ======
        if os.getenv("GOOGLE_API_KEY"):
            llm = ChatGoogleGenerativeAI(
                model="gemini-2.5-flash",
                google_api_key=os.getenv("GOOGLE_API_KEY")
            )
        else:
            st.warning("⚠️ Tidak ada API Key, chatbot nonaktif.")
            llm = None

        if llm:
            tools = [PythonREPLTool()]
            agent = initialize_agent(
                tools=tools,
                llm=llm,
                agent_type=AgentType.ZERO_SHOT_REACT_DESCRIPTION,
                verbose=True
            )

            user_query = st.chat_input("Tanyakan sesuatu tentang data (statistik, tren, dsb.)")
            if user_query:
                with st.chat_message("user"):
                    st.markdown(user_query)
                with st.spinner("🔎 Menganalisis data..."):
                    try:
                        # Kirim pertanyaan ke LLM dengan konteks data (CSV string ringkas)
                        preview = df.head(50).to_csv(index=False)
                        prompt = f"Gunakan Python untuk analisa data berikut:\n{preview}\n\nPertanyaan: {user_query}"
                        response = agent.run(prompt)
                        with st.chat_message("assistant"):
                            st.markdown(response)
                    except Exception as e:
                        st.error(f"❌ Error: {e}")

# ========== MODE 2: RAG Advanced ==========
if menu == "📑 RAG Advanced":
    uploaded_files = st.file_uploader(
        "Upload dokumen (PDF, TXT, DOCX, PPTX, CSV, XLSX, gambar dengan teks) → bisa multi-file",
        type=["pdf", "txt", "docx", "pptx", "csv", "xls", "xlsx", "png", "jpg", "jpeg", "bmp"],
        accept_multiple_files=True
    )

    if uploaded_files:
        st.markdown("### 📂 Dokumen yang diupload:")
        for f in uploaded_files:
            st.write("- " + f.name)

        vectorstore = build_vectorstore(uploaded_files)
        retriever = vectorstore.as_retriever(search_kwargs={"k": 3})

        if "qa_chain" not in st.session_state:
            if os.getenv("GOOGLE_API_KEY"):
                llm = ChatGoogleGenerativeAI(
                    model="gemini-2.5-flash",
                    google_api_key=os.getenv("GOOGLE_API_KEY")
                )
            else:
                st.stop()

            st.session_state.qa_chain = RetrievalQA.from_chain_type(
                llm=llm,
                retriever=retriever,
                return_source_documents=True
            )
            st.success("✅ Chatbot RAG siap digunakan")

        user_query = st.chat_input("Tanyakan sesuatu tentang dokumen hukum/medis...")
        if user_query:
            with st.chat_message("user"):
                st.markdown(user_query)

            with st.spinner("🔎 Menganalisis dokumen..."):
                try:
                    result = st.session_state.qa_chain({"query": user_query})
                    answer = result["result"]
                    sources = result.get("source_documents", [])

                    with st.chat_message("assistant"):
                        st.markdown(answer)
                        if sources:
                            st.write("**Sumber:**")
                            for s in sources:
                                st.caption(f"{s.metadata.get('source','')} → {s.page_content[:200]}...")
                except Exception as e:
                    st.error(f"❌ Error: {e}")
